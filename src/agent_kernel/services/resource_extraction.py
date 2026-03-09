"""Resource Extraction Service - Extract and summarize linked resources (v1.0.6).

Handles extraction and summarization of various file types linked in notes:
- PowerPoint (.pptx, .ppt)
- PDF documents (.pdf)
- Word documents (.docx, .doc)
- Excel spreadsheets (.xlsx, .xls)
- Text/markdown files

Creates summary notes in Obsidian organized by project or area.
"""

from __future__ import annotations

import hashlib
import mimetypes
import re
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import TYPE_CHECKING, Any

import structlog

from agent_kernel.core.schemas.base import utc_now

if TYPE_CHECKING:
    from agent_kernel.services.llm import LLMService

logger = structlog.get_logger(__name__)


class ResourceType(str, Enum):
    """Supported resource types for extraction."""

    POWERPOINT = "powerpoint"
    PDF = "pdf"
    WORD = "word"
    EXCEL = "excel"
    TEXT = "text"
    MARKDOWN = "markdown"
    UNKNOWN = "unknown"


# File extension to resource type mapping
EXTENSION_MAPPING = {
    ".pptx": ResourceType.POWERPOINT,
    ".ppt": ResourceType.POWERPOINT,
    ".pdf": ResourceType.PDF,
    ".docx": ResourceType.WORD,
    ".doc": ResourceType.WORD,
    ".xlsx": ResourceType.EXCEL,
    ".xls": ResourceType.EXCEL,
    ".txt": ResourceType.TEXT,
    ".md": ResourceType.MARKDOWN,
}


@dataclass
class ResourceMetadata:
    """Metadata about an extracted resource."""

    file_path: str
    file_name: str
    resource_type: ResourceType
    file_size: int = 0
    created_at: datetime | None = None
    modified_at: datetime | None = None
    content_hash: str = ""
    page_count: int | None = None
    slide_count: int | None = None
    word_count: int | None = None

    def to_dict(self) -> dict[str, Any]:
        return {
            "file_path": self.file_path,
            "file_name": self.file_name,
            "resource_type": self.resource_type.value,
            "file_size": self.file_size,
            "created_at": self.created_at.isoformat() if self.created_at else None,
            "modified_at": self.modified_at.isoformat() if self.modified_at else None,
            "content_hash": self.content_hash,
            "page_count": self.page_count,
            "slide_count": self.slide_count,
            "word_count": self.word_count,
        }


@dataclass
class ExtractionResult:
    """Result of extracting content from a resource."""

    metadata: ResourceMetadata
    raw_text: str = ""
    structured_content: dict[str, Any] = field(default_factory=dict)
    extraction_method: str = ""
    success: bool = True
    error: str | None = None

    def to_dict(self) -> dict[str, Any]:
        return {
            "metadata": self.metadata.to_dict(),
            "raw_text_length": len(self.raw_text),
            "extraction_method": self.extraction_method,
            "success": self.success,
            "error": self.error,
        }


@dataclass
class ResourceSummary:
    """LLM-generated summary of a resource."""

    title: str = ""
    summary: str = ""
    key_points: list[str] = field(default_factory=list)
    topics: list[str] = field(default_factory=list)
    action_items: list[str] = field(default_factory=list)
    suggested_project: str | None = None
    confidence: float = 0.0

    def to_dict(self) -> dict[str, Any]:
        return {
            "title": self.title,
            "summary": self.summary,
            "key_points": self.key_points,
            "topics": self.topics,
            "action_items": self.action_items,
            "suggested_project": self.suggested_project,
            "confidence": self.confidence,
        }

    def to_markdown(self, file_link: str, metadata: ResourceMetadata) -> str:
        """Generate markdown content for a summary note."""
        lines = [
            f"# {self.title}",
            "",
            f"**Source:** {file_link}",
            f"**Type:** {metadata.resource_type.value.title()}",
        ]

        if metadata.page_count:
            lines.append(f"**Pages:** {metadata.page_count}")
        if metadata.slide_count:
            lines.append(f"**Slides:** {metadata.slide_count}")

        lines.extend(["", "## Summary", "", self.summary, ""])

        if self.key_points:
            lines.append("## Key Points")
            lines.append("")
            for point in self.key_points:
                lines.append(f"- {point}")
            lines.append("")

        if self.action_items:
            lines.append("## Action Items")
            lines.append("")
            for item in self.action_items:
                lines.append(f"- [ ] {item}")
            lines.append("")

        if self.topics:
            lines.append("## Topics")
            lines.append("")
            lines.append(", ".join([f"`{t}`" for t in self.topics]))
            lines.append("")

        return "\n".join(lines)


# Summarization prompt templates
RESOURCE_SUMMARY_SYSTEM_PROMPT = """\
You are an expert at analyzing documents and creating concise, actionable summaries.

Your task is to analyze the extracted content and produce:
1. **Title**: A clear, descriptive title for the document
2. **Summary**: A 2-3 paragraph executive summary
3. **Key Points**: 5-10 bullet points of the most important information
4. **Topics**: 3-5 topic tags (lowercase, hyphenated)
5. **Action Items**: Any tasks or follow-ups mentioned (optional)
6. **Project Suggestion**: Which project/area this belongs to

Consider the content type ({resource_type}) when structuring the summary:
- For presentations: Focus on main messages and conclusions
- For documents: Focus on key findings, decisions, and recommendations
- For spreadsheets: Focus on data patterns and insights

AVAILABLE PROJECTS: {projects}

Respond in JSON format:
{{
  "title": "Document Title",
  "summary": "Executive summary paragraph(s)...",
  "key_points": ["Point 1", "Point 2", ...],
  "topics": ["topic-1", "topic-2"],
  "action_items": ["Action 1", "Action 2"],
  "suggested_project": "project_name",
  "confidence": 0.85
}}
"""

RESOURCE_SUMMARY_USER_PROMPT = """\
Summarize this {resource_type} document:

---
File: {file_name}
{metadata_section}
---

Content:
{content}

---

Respond with JSON only, no markdown formatting.
"""


class ResourceExtractor:
    """Extracts text content from various file types."""

    def __init__(self) -> None:
        """Initialize the resource extractor."""
        self._extractors: dict[ResourceType, callable] = {
            ResourceType.POWERPOINT: self._extract_powerpoint,
            ResourceType.PDF: self._extract_pdf,
            ResourceType.WORD: self._extract_word,
            ResourceType.EXCEL: self._extract_excel,
            ResourceType.TEXT: self._extract_text,
            ResourceType.MARKDOWN: self._extract_text,
        }

    def get_resource_type(self, file_path: str | Path) -> ResourceType:
        """Determine resource type from file path."""
        path = Path(file_path)
        ext = path.suffix.lower()
        return EXTENSION_MAPPING.get(ext, ResourceType.UNKNOWN)

    def _compute_content_hash(self, file_path: Path) -> str:
        """Compute hash of file content."""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()[:16]

    def _get_file_metadata(self, file_path: Path) -> ResourceMetadata:
        """Get basic file metadata."""
        stat = file_path.stat()
        return ResourceMetadata(
            file_path=str(file_path),
            file_name=file_path.name,
            resource_type=self.get_resource_type(file_path),
            file_size=stat.st_size,
            created_at=datetime.fromtimestamp(stat.st_ctime),
            modified_at=datetime.fromtimestamp(stat.st_mtime),
            content_hash=self._compute_content_hash(file_path),
        )

    async def extract(self, file_path: str | Path) -> ExtractionResult:
        """Extract content from a file.

        Args:
            file_path: Path to the file to extract.

        Returns:
            ExtractionResult with extracted content.
        """
        path = Path(file_path)

        if not path.exists():
            return ExtractionResult(
                metadata=ResourceMetadata(
                    file_path=str(path),
                    file_name=path.name,
                    resource_type=ResourceType.UNKNOWN,
                ),
                success=False,
                error=f"File not found: {path}",
            )

        metadata = self._get_file_metadata(path)
        resource_type = metadata.resource_type

        if resource_type == ResourceType.UNKNOWN:
            return ExtractionResult(
                metadata=metadata,
                success=False,
                error=f"Unsupported file type: {path.suffix}",
            )

        extractor = self._extractors.get(resource_type)
        if not extractor:
            return ExtractionResult(
                metadata=metadata,
                success=False,
                error=f"No extractor for type: {resource_type.value}",
            )

        try:
            result = await extractor(path, metadata)
            return result
        except Exception as e:
            logger.exception("extraction_failed", file_path=str(path))
            return ExtractionResult(
                metadata=metadata,
                success=False,
                error=str(e),
            )

    async def _extract_powerpoint(
        self,
        file_path: Path,
        metadata: ResourceMetadata,
    ) -> ExtractionResult:
        """Extract content from PowerPoint files."""
        try:
            from pptx import Presentation
        except ImportError:
            return ExtractionResult(
                metadata=metadata,
                success=False,
                error="python-pptx not installed. Install with: pip install python-pptx",
            )

        prs = Presentation(str(file_path))
        slides_content: list[dict[str, Any]] = []
        all_text: list[str] = []

        for idx, slide in enumerate(prs.slides, 1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            slides_content.append({
                "slide_number": idx,
                "text": "\n".join(slide_text),
            })
            all_text.extend(slide_text)

        metadata.slide_count = len(prs.slides)
        metadata.word_count = len(" ".join(all_text).split())

        return ExtractionResult(
            metadata=metadata,
            raw_text="\n\n".join(all_text),
            structured_content={"slides": slides_content},
            extraction_method="python-pptx",
            success=True,
        )

    async def _extract_pdf(
        self,
        file_path: Path,
        metadata: ResourceMetadata,
    ) -> ExtractionResult:
        """Extract content from PDF files."""
        try:
            import pypdf
        except ImportError:
            return ExtractionResult(
                metadata=metadata,
                success=False,
                error="pypdf not installed. Install with: pip install pypdf",
            )

        reader = pypdf.PdfReader(str(file_path))
        pages_content: list[dict[str, Any]] = []
        all_text: list[str] = []

        for idx, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            pages_content.append({
                "page_number": idx,
                "text": text,
            })
            all_text.append(text)

        metadata.page_count = len(reader.pages)
        metadata.word_count = len(" ".join(all_text).split())

        return ExtractionResult(
            metadata=metadata,
            raw_text="\n\n".join(all_text),
            structured_content={"pages": pages_content},
            extraction_method="pypdf",
            success=True,
        )

    async def _extract_word(
        self,
        file_path: Path,
        metadata: ResourceMetadata,
    ) -> ExtractionResult:
        """Extract content from Word documents."""
        try:
            import docx
        except ImportError:
            return ExtractionResult(
                metadata=metadata,
                success=False,
                error="python-docx not installed. Install with: pip install python-docx",
            )

        doc = docx.Document(str(file_path))
        paragraphs: list[str] = []
        tables: list[list[list[str]]] = []

        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        for table in doc.tables:
            table_data: list[list[str]] = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)

        all_text = "\n\n".join(paragraphs)
        metadata.word_count = len(all_text.split())

        return ExtractionResult(
            metadata=metadata,
            raw_text=all_text,
            structured_content={
                "paragraphs": paragraphs,
                "tables": tables,
            },
            extraction_method="python-docx",
            success=True,
        )

    async def _extract_excel(
        self,
        file_path: Path,
        metadata: ResourceMetadata,
    ) -> ExtractionResult:
        """Extract content from Excel spreadsheets."""
        try:
            import openpyxl
        except ImportError:
            return ExtractionResult(
                metadata=metadata,
                success=False,
                error="openpyxl not installed. Install with: pip install openpyxl",
            )

        wb = openpyxl.load_workbook(str(file_path), data_only=True)
        sheets_content: dict[str, list[list[Any]]] = {}
        all_text: list[str] = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows: list[list[Any]] = []
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                rows.append(row_values)
                all_text.extend([v for v in row_values if v])
            sheets_content[sheet_name] = rows

        metadata.word_count = len(" ".join(all_text).split())

        return ExtractionResult(
            metadata=metadata,
            raw_text="\n".join(all_text),
            structured_content={"sheets": sheets_content},
            extraction_method="openpyxl",
            success=True,
        )

    async def _extract_text(
        self,
        file_path: Path,
        metadata: ResourceMetadata,
    ) -> ExtractionResult:
        """Extract content from text/markdown files."""
        try:
            text = file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            text = file_path.read_text(encoding="latin-1")

        metadata.word_count = len(text.split())

        return ExtractionResult(
            metadata=metadata,
            raw_text=text,
            extraction_method="text",
            success=True,
        )


class ResourceExtractionService:
    """Service for extracting and summarizing linked resources.

    Provides:
    - Content extraction from various file types
    - LLM-powered summarization
    - Note generation for Obsidian
    - Project-based organization
    """

    def __init__(
        self,
        llm_service: LLMService,
        vault_path: str | Path | None = None,
        extractor: ResourceExtractor | None = None,
        available_projects: list[str] | None = None,
    ) -> None:
        """Initialize the resource extraction service.

        Args:
            llm_service: LLM service for summarization.
            vault_path: Path to Obsidian vault for note generation.
            extractor: Custom extractor (uses default if None).
            available_projects: List of available project names.
        """
        self.llm_service = llm_service
        self.vault_path = Path(vault_path) if vault_path else None
        self.extractor = extractor or ResourceExtractor()
        self.available_projects = available_projects or []

        # Track processed resources for idempotency
        self._processed: dict[str, str] = {}  # content_hash -> note_path

        logger.info(
            "resource_extraction_service_initialized",
            vault_path=str(self.vault_path) if self.vault_path else None,
            projects=len(self.available_projects),
        )

    def find_resource_links(self, markdown_content: str) -> list[str]:
        """Find resource file links in markdown content.

        Args:
            markdown_content: Markdown text to search.

        Returns:
            List of file paths found.
        """
        # Match wiki-style links: [[file.pptx]]
        wiki_pattern = r"\[\[([^\]]+\.(?:pptx?|pdf|docx?|xlsx?|txt))\]\]"
        wiki_matches = re.findall(wiki_pattern, markdown_content, re.IGNORECASE)

        # Match markdown links: [text](file.pptx)
        md_pattern = r"\[([^\]]*)\]\(([^)]+\.(?:pptx?|pdf|docx?|xlsx?|txt))\)"
        md_matches = [m[1] for m in re.findall(md_pattern, markdown_content, re.IGNORECASE)]

        # Match bare file paths
        bare_pattern = r"(?:^|\s)([./\w-]+\.(?:pptx?|pdf|docx?|xlsx?|txt))(?:\s|$)"
        bare_matches = re.findall(bare_pattern, markdown_content, re.IGNORECASE)

        all_links = list(set(wiki_matches + md_matches + bare_matches))
        return all_links

    async def extract_and_summarize(
        self,
        file_path: str | Path,
        max_content_length: int = 8000,
    ) -> tuple[ExtractionResult, ResourceSummary | None]:
        """Extract content and generate summary for a resource.

        Args:
            file_path: Path to the resource file.
            max_content_length: Max content chars to send to LLM.

        Returns:
            Tuple of (ExtractionResult, ResourceSummary or None if failed).
        """
        # Extract content
        extraction = await self.extractor.extract(file_path)

        if not extraction.success:
            return extraction, None

        # Check if already processed
        if extraction.metadata.content_hash in self._processed:
            logger.debug(
                "resource_already_processed",
                file_path=str(file_path),
                hash=extraction.metadata.content_hash,
            )

        # Generate summary
        summary = await self._generate_summary(extraction, max_content_length)

        return extraction, summary

    async def _generate_summary(
        self,
        extraction: ExtractionResult,
        max_content_length: int,
    ) -> ResourceSummary:
        """Generate LLM summary of extracted content."""
        # Truncate content if needed
        content = extraction.raw_text[:max_content_length]
        if len(extraction.raw_text) > max_content_length:
            content += "\n\n[Content truncated...]"

        # Build metadata section
        metadata_lines = []
        if extraction.metadata.slide_count:
            metadata_lines.append(f"Slides: {extraction.metadata.slide_count}")
        if extraction.metadata.page_count:
            metadata_lines.append(f"Pages: {extraction.metadata.page_count}")
        if extraction.metadata.word_count:
            metadata_lines.append(f"Words: ~{extraction.metadata.word_count}")

        metadata_section = "\n".join(metadata_lines) if metadata_lines else ""

        # Build prompts
        system_prompt = RESOURCE_SUMMARY_SYSTEM_PROMPT.format(
            resource_type=extraction.metadata.resource_type.value,
            projects=", ".join(self.available_projects) if self.available_projects else "None",
        )

        user_prompt = RESOURCE_SUMMARY_USER_PROMPT.format(
            resource_type=extraction.metadata.resource_type.value,
            file_name=extraction.metadata.file_name,
            metadata_section=metadata_section,
            content=content,
        )

        # Call LLM
        response = await self.llm_service.generate(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0.3,
            max_tokens=1000,
        )

        return self._parse_summary(response)

    def _parse_summary(self, response: str) -> ResourceSummary:
        """Parse LLM response into ResourceSummary."""
        summary = ResourceSummary()

        # Clean response
        text = response.strip()
        if text.startswith("```"):
            first_newline = text.find("\n")
            if first_newline > 0:
                text = text[first_newline + 1:]
            text = text.removesuffix("```").strip()

        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            logger.warning("failed_to_parse_summary_response", response=response[:200])
            return summary

        summary.title = data.get("title", "")
        summary.summary = data.get("summary", "")
        summary.key_points = data.get("key_points", [])
        summary.topics = data.get("topics", [])
        summary.action_items = data.get("action_items", [])
        summary.suggested_project = data.get("suggested_project")
        summary.confidence = float(data.get("confidence", 0.0))

        return summary

    async def create_summary_note(
        self,
        extraction: ExtractionResult,
        summary: ResourceSummary,
        target_folder: str | None = None,
        note_name: str | None = None,
    ) -> str | None:
        """Create an Obsidian note with the resource summary.

        Args:
            extraction: The extraction result.
            summary: The generated summary.
            target_folder: Target folder in vault (uses project suggestion if None).
            note_name: Custom note name (derives from summary title if None).

        Returns:
            Path to created note, or None if vault_path not set.
        """
        if not self.vault_path:
            logger.warning("vault_path_not_set_cannot_create_note")
            return None

        # Determine target folder
        if target_folder is None:
            if summary.suggested_project:
                target_folder = f"Resources/{summary.suggested_project}"
            else:
                target_folder = "Resources/Unsorted"

        # Determine note name
        if note_name is None:
            # Sanitize title for filename
            safe_title = re.sub(r'[<>:"/\\|?*]', '', summary.title or "Untitled")
            safe_title = safe_title.strip()[:100]  # Limit length
            note_name = f"{safe_title}.md"

        # Create folder if needed
        folder_path = self.vault_path / target_folder
        folder_path.mkdir(parents=True, exist_ok=True)

        # Generate file link (relative to vault)
        file_link = f"[[{extraction.metadata.file_name}]]"

        # Generate markdown content
        content = summary.to_markdown(file_link, extraction.metadata)

        # Add frontmatter
        frontmatter = [
            "---",
            f"source_file: \"{extraction.metadata.file_path}\"",
            f"source_type: {extraction.metadata.resource_type.value}",
            f"content_hash: {extraction.metadata.content_hash}",
            f"extracted_at: {utc_now().isoformat()}",
        ]
        if summary.topics:
            frontmatter.append(f"tags: [{', '.join(summary.topics)}]")
        if summary.suggested_project:
            frontmatter.append(f"project: {summary.suggested_project}")
        frontmatter.append("auto:")
        frontmatter.append(f"  generated: true")
        frontmatter.append(f"  confidence: {summary.confidence}")
        frontmatter.append("---")
        frontmatter.append("")

        full_content = "\n".join(frontmatter) + content

        # Write note
        note_path = folder_path / note_name
        note_path.write_text(full_content, encoding="utf-8")

        # Track as processed
        self._processed[extraction.metadata.content_hash] = str(note_path)

        logger.info(
            "summary_note_created",
            note_path=str(note_path),
            source_file=extraction.metadata.file_name,
        )

        return str(note_path)

    async def process_note_resources(
        self,
        note_path: str | Path,
        create_notes: bool = True,
    ) -> list[tuple[ExtractionResult, ResourceSummary | None, str | None]]:
        """Process all resource links in a note.

        Args:
            note_path: Path to the note to process.
            create_notes: Whether to create summary notes.

        Returns:
            List of (ExtractionResult, ResourceSummary, note_path) tuples.
        """
        note = Path(note_path)
        if not note.exists():
            logger.warning("note_not_found", path=str(note))
            return []

        content = note.read_text(encoding="utf-8")
        links = self.find_resource_links(content)

        if not links:
            return []

        results: list[tuple[ExtractionResult, ResourceSummary | None, str | None]] = []

        for link in links:
            # Resolve relative paths
            if not Path(link).is_absolute():
                resource_path = note.parent / link
            else:
                resource_path = Path(link)

            extraction, summary = await self.extract_and_summarize(resource_path)

            note_created = None
            if create_notes and extraction.success and summary:
                note_created = await self.create_summary_note(extraction, summary)

            results.append((extraction, summary, note_created))

        return results

    def is_processed(self, file_path: str | Path) -> bool:
        """Check if a resource has already been processed."""
        path = Path(file_path)
        if not path.exists():
            return False
        hash_val = self.extractor._compute_content_hash(path)
        return hash_val in self._processed

    def get_processed_note(self, file_path: str | Path) -> str | None:
        """Get the note path for a processed resource."""
        path = Path(file_path)
        if not path.exists():
            return None
        hash_val = self.extractor._compute_content_hash(path)
        return self._processed.get(hash_val)


# Import json at module level for use in parsing
import json
